"""PPTX file parser using python-pptx (MIT license).

Extracts text from text frames, tables, **and** embedded images (via OCR)
on each slide.  Shapes are processed in slide order so that image OCR
text appears at the correct position relative to surrounding content.
"""

from __future__ import annotations

import io
import logging
from pathlib import Path

from rag_kb.parsers.base import ParsedDocument

logger = logging.getLogger(__name__)


class PptxParser:
    """Parse Microsoft PowerPoint .pptx files."""

    supported_extensions: list[str] = [".pptx"]

    def parse(self, file_path: Path) -> ParsedDocument:
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError as exc:
            raise ImportError(
                "python-pptx is required for PPTX parsing: pip install python-pptx"
            ) from exc

        prs = Presentation(str(file_path))
        slides_text: list[str] = []
        total_images_ocrd = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts: list[str] = [f"[Slide {slide_num}]"]
            for shape in slide.shapes:
                # --- Text frame ---
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_parts.append(text)

                # --- Table ---
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_parts.append(row_text)

                # --- Picture / image shape → OCR ---
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    ocr_text = self._ocr_shape_image(shape, slide_num)
                    if ocr_text:
                        slide_parts.append(f"[Image OCR]\n{ocr_text}")
                        total_images_ocrd += 1

            if len(slide_parts) > 1:  # more than just the slide marker
                slides_text.append("\n".join(slide_parts))

        text = "\n\n".join(slides_text)

        metadata: dict[str, str] = {
            "format": "pptx",
            "slide_count": str(len(prs.slides)),
        }
        if total_images_ocrd:
            metadata["images_ocrd"] = str(total_images_ocrd)

        return ParsedDocument(
            text=text,
            source_path=str(file_path),
            metadata=metadata,
        )

    # ------------------------------------------------------------------

    @staticmethod
    def _ocr_shape_image(shape, slide_num: int) -> str:
        """OCR a picture shape and return the extracted text (or "")."""
        from PIL import Image as PILImage

        from rag_kb.parsers.image_parser import ocr_image

        try:
            blob = shape.image.blob
            pil_img = PILImage.open(io.BytesIO(blob)).convert("RGB")

            # Skip tiny images (icons, logos, decorative)
            if pil_img.width < 50 or pil_img.height < 50:
                return ""

            label = f"slide-{slide_num}/{shape.name}"
            ocr_text, _engine = ocr_image(pil_img, label=label)
            return ocr_text

        except Exception as exc:
            logger.debug(
                "PPTX image OCR failed for shape '%s' on slide %d: %s", shape.name, slide_num, exc
            )
            return ""
